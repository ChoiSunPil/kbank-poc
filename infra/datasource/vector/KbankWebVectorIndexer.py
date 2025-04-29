# web_faiss_chroma.py
import time
from typing import List

from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from bs4 import BeautifulSoup
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from collections import deque
import re
from dotenv import load_dotenv
import os
import pandas as pd
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation

class KbankWebVectorIndexer:
    def __init__(self):
        load_dotenv()
        self.persist_dir = os.getenv("PERSIST_DIR")
        self.embeddings = HuggingFaceEmbeddings(
            model_name="jhgan/ko-sbert-sts",
            model_kwargs={"device": "cpu"}
        )
        if self._has_files(self.persist_dir):
            self.vector_db = Chroma(
                embedding_function=self.embeddings,
                persist_directory=self.persist_dir
            )
            print("✅ 기존 벡터 DB 로드 완료")
        else:
            self.vector_db = None
            print("📁 벡터 DB 없음 — 새로 생성 필요")
        self.documents: list[Document] = []
        self.visited = set()
        self.max_pages = 30  # 크롤링할 최대 페이지 수

    def _has_files(self, dir_path: str) -> bool:
        return os.path.isdir(dir_path) and len(os.listdir(dir_path)) > 0
    def crawl(self, start_path: str):
        print("🌐 하위 페이지 포함 전체 크롤링 중...")
        queue = deque([self.base_url + start_path])
        count = 0

        while queue and count < self.max_pages:
            time.sleep(0.3)
            current_url = queue.popleft()
            if current_url in self.visited:
                continue

            print(current_url)

            try:
                self.visited.add(current_url)
                count += 1
                service = Service("/usr/local/bin/chromedriver")
                driver = webdriver.Chrome(service=service)
                driver.get(current_url)
                html = driver.page_source
                soup = BeautifulSoup(html, "html.parser")
                text = soup.get_text(separator="\n", strip=True)

                # 텍스트가 충분히 길면만 벡터화
                if len(text) > 100:
                    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
                    chunks = splitter.split_text(text)
                    self.documents.extend([
                        Document(page_content=chunk, metadata={"source": current_url})
                        for chunk in chunks
                    ])

                # 링크 추출 및 큐에 추가
                pattern = re.compile(r"goMenu\(['\"]([^'\"]+)['\"]\)")
                matches = pattern.findall(html)
                for menu_id in matches:
                    menu_url = self.base_url + f"/ib20/mnu/{menu_id}"
                    if menu_url not in self.visited:
                        queue.append(menu_url)

                print(f"✅ 수집 완료: {current_url}")

            except Exception as e:
                print(f"⚠️ 실패: {current_url} | 이유: {e}")

    def build_vector_db(self):
        self.vector_db = Chroma.from_documents(
            documents=self.documents,
            embedding=self.embeddings,
            persist_directory=self.persist_dir
        )

    def search(self, query: str, top_k: int = 3) -> list:
        if not self.vector_db:
            raise RuntimeError("❗ 벡터 DB가 준비되지 않았습니다. 먼저 build_vector_db()를 호출하세요.")
        print(f"🔍 검색 중: '{query}'")
        results = self.vector_db.similarity_search(query, k=top_k)
        for i, doc in enumerate(results):
            print(f"\n🔹 결과 {i + 1}")
            print(doc.page_content)
            print(f"[출처]: {doc.metadata.get('source')}")
        return results

    def file_to_vector(self,file):
        extracted_text = self.__extract_text(file)
        chunks = self.__split_text(extracted_text, file.name)
        self.__update_vector_db(chunks)
        return None

    def __extract_text(self,file):
        ext = Path(file.name).suffix.lower()
        if ext == ".pdf":
            reader = PdfReader(file)
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext in [".ppt", ".pptx"]:
            prs = Presentation(file)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext in [".csv", ".tsv"]:
            df = pd.read_csv(file)
            return df.to_string()
        elif ext in [".xlsx", ".xls"]:
            df = pd.read_excel(file)
            return df.to_string()
        else:
            raise ValueError("지원하지 않는 파일 형식입니다.")

    def __split_text(self, text:str, source:str) -> List[Document]:
        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        chunks = splitter.split_text(text)
        return [Document(page_content=chunk, metadata={"source":source}) for chunk in chunks]

    def __update_vector_db(self, documents: list[Document]):
        if not self.vector_db:
            self.vector_db = Chroma.from_documents(documents, embedding=self.embeddings, persist_directory=self.persist_dir)
        else:
            self.vector_db.add_documents(documents)
